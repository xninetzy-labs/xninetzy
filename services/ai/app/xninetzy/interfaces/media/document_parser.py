"""Parse documents into plain text.

Each format is handled best-effort. Optional libraries (python-docx, openpyxl,
python-pptx) are imported lazily so the module loads even when they are not
installed; if a needed library is missing, ``parse_document`` returns a clear
error instead of raising.

PDF reuses the existing pypdf-based reader (``tools.hebat.pdf_reader``).
"""
from __future__ import annotations

import json
from pathlib import Path

from app.xninetzy.core.logging import logging

logger = logging.getLogger(__name__)

# Extensions we can turn into text in this slice.
SUPPORTED_DOC_EXTS = {
    ".pdf", ".txt", ".md", ".markdown", ".csv", ".json",
    ".docx", ".xlsx", ".pptx",
}

# mime -> extension hint, used when filename has no usable extension.
_MIME_EXT = {
    "application/pdf": ".pdf",
    "text/plain": ".txt",
    "text/markdown": ".md",
    "text/csv": ".csv",
    "application/json": ".json",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
}


def _resolve_ext(path: Path, mime_type: str | None, filename: str | None) -> str:
    for candidate in (path.suffix, Path(filename or "").suffix):
        if candidate and candidate.lower() in SUPPORTED_DOC_EXTS:
            return candidate.lower()
    if mime_type:
        return _MIME_EXT.get(mime_type.split(";")[0].strip().lower(), "")
    return ""


def _err(msg: str) -> dict:
    return {"text": "", "char_count": 0, "kind": None, "error": msg}


def parse_document(path: str, mime_type: str | None = None, filename: str | None = None) -> dict:
    """Parse a document file into text.

    Returns ``{text, char_count, kind, error}``. ``error`` is None on success.
    """
    p = Path(path)
    if not p.exists() or not p.is_file():
        return _err(f"File tidak ditemukan: {filename or p.name}")

    ext = _resolve_ext(p, mime_type, filename)
    if not ext:
        return _err(f"Tipe file belum didukung (mime={mime_type}, name={filename or p.name})")

    try:
        if ext == ".pdf":
            text = _parse_pdf(p)
        elif ext in (".txt", ".md", ".markdown", ".csv", ".json"):
            text = _parse_text(p)
        elif ext == ".docx":
            text = _parse_docx(p)
        elif ext == ".xlsx":
            text = _parse_xlsx(p)
        elif ext == ".pptx":
            text = _parse_pptx(p)
        else:
            return _err(f"Ekstensi {ext} belum didukung")
    except _MissingLib as exc:
        return _err(str(exc))
    except Exception as exc:  # pragma: no cover - defensive
        logger.warning("Document parse failed for %s: %s", p.name, exc)
        return _err(f"Gagal membaca dokumen: {exc}")

    text = (text or "").strip()
    if not text:
        return _err("Dokumen terbaca tetapi tidak ada teks yang bisa diekstrak.")
    return {"text": text, "char_count": len(text), "kind": ext.lstrip("."), "error": None}


class _MissingLib(RuntimeError):
    pass


def _parse_pdf(p: Path) -> str:
    from app.xninetzy.os.academic.hebat.pdf_reader import read_pdf_text

    result = read_pdf_text(str(p))
    if result.get("error"):
        raise RuntimeError(result["error"])
    return result.get("text", "")


def _parse_text(p: Path) -> str:
    raw = p.read_text(encoding="utf-8", errors="replace")
    if p.suffix.lower() == ".json":
        try:
            return json.dumps(json.loads(raw), ensure_ascii=False, indent=2)
        except Exception:
            return raw
    return raw


def _parse_docx(p: Path) -> str:
    try:
        import docx  # python-docx
    except ImportError as exc:
        raise _MissingLib("Library python-docx belum terinstall untuk membaca .docx") from exc
    doc = docx.Document(str(p))
    parts = [para.text for para in doc.paragraphs if para.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _parse_xlsx(p: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise _MissingLib("Library openpyxl belum terinstall untuk membaca .xlsx") from exc
    wb = load_workbook(str(p), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _parse_pptx(p: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise _MissingLib("Library python-pptx belum terinstall untuk membaca .pptx") from exc
    prs = Presentation(str(p))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
    return "\n".join(parts)
