"""Office extractor — DOCX/PPTX/XLSX into structural blocks.

Reuses the same lazy-import discipline as ``document_parser``: a missing
optional lib raises ``_MissingLib`` with a clear message rather than a bare
ImportError.
"""
from __future__ import annotations

from pathlib import Path

from app.xninetzy.interfaces.media.document_parser import _MissingLib, _resolve_ext
from app.xninetzy.os.knowledge.extraction.extractors.tables import rows_to_markdown
from app.xninetzy.os.knowledge.extraction.schemas import DocBlock, StructuredDocument


def extract_office(path: str, extract_tables: bool = True) -> StructuredDocument:
    ext = _resolve_ext(Path(path), None, None)
    if ext == ".docx":
        return _docx(path, extract_tables)
    if ext == ".pptx":
        return _pptx(path)
    if ext == ".xlsx":
        return _xlsx(path)
    raise _MissingLib(f"Ekstensi office {ext} belum didukung extractor")


def _docx(path: str, extract_tables: bool) -> StructuredDocument:
    try:
        import docx  # python-docx
    except ImportError as exc:
        raise _MissingLib("Library python-docx belum terinstall untuk membaca .docx") from exc

    doc = docx.Document(str(path))
    blocks: list[DocBlock] = []
    heading_stack: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name if para.style else "") or ""
        if style.startswith("Heading") or style == "Title":
            digits = "".join(ch for ch in style if ch.isdigit())
            level = int(digits) if digits else 1
            # Truncate deeper levels and set this level's heading.
            heading_stack[level - 1:] = [text]
            blocks.append(
                DocBlock(kind="heading", text=text, heading_path=tuple(heading_stack))
            )
        else:
            kind = "list" if para.style and "List" in (para.style.name or "") else "paragraph"
            blocks.append(
                DocBlock(kind=kind, text=text, heading_path=tuple(heading_stack))
            )

    table_count = 0
    if extract_tables:
        for table in doc.tables:
            rows = [[c.text for c in row.cells] for row in table.rows]
            md = rows_to_markdown(rows)
            if md:
                table_count += 1
                blocks.append(
                    DocBlock(kind="table", text=md, heading_path=tuple(heading_stack),
                             meta={"rows": len(rows)})
                )

    return StructuredDocument(
        blocks=tuple(blocks), page_count=1, table_count=table_count, kind="docx"
    )


def _pptx(path: str) -> StructuredDocument:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise _MissingLib("Library python-pptx belum terinstall untuk membaca .pptx") from exc

    prs = Presentation(str(path))
    blocks: list[DocBlock] = []
    slide_count = 0
    for i, slide in enumerate(prs.slides, 1):
        slide_count = i
        title = f"Slide {i}"
        heading_path = (title,)
        blocks.append(DocBlock(kind="heading", text=title, page=i, heading_path=heading_path))
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    blocks.append(
                        DocBlock(kind="paragraph", text=line, page=i, heading_path=heading_path)
                    )
    return StructuredDocument(
        blocks=tuple(blocks), page_count=slide_count, kind="pptx"
    )


def _xlsx(path: str) -> StructuredDocument:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise _MissingLib("Library openpyxl belum terinstall untuk membaca .xlsx") from exc

    wb = load_workbook(str(path), read_only=True, data_only=True)
    blocks: list[DocBlock] = []
    table_count = 0
    for ws in wb.worksheets:
        heading_path = (ws.title,)
        blocks.append(DocBlock(kind="heading", text=ws.title, heading_path=heading_path))
        rows = [
            [c for c in row]
            for row in ws.iter_rows(values_only=True)
            if any(c is not None for c in row)
        ]
        md = rows_to_markdown(rows)
        if md:
            table_count += 1
            blocks.append(
                DocBlock(kind="table", text=md, heading_path=heading_path,
                         meta={"sheet": ws.title, "rows": len(rows)})
            )
    return StructuredDocument(
        blocks=tuple(blocks), page_count=len(wb.worksheets), table_count=table_count,
        kind="xlsx",
    )
